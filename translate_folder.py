
from googletrans import Translator, constants
from pprint import pprint
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
import pdb
import os
import sys

#change this to target language
LANG='es'
#change this to source folder
SOURCE = 'MOOC2'
#change this to result folder
RESULT = 'MOOC2_T'
#TODO-make this callable with parameters.

def replaceText(_paragraph, _replaceString):
        paragraph = _paragraph
        replaced = None
        for rns in paragraph.runs:
            if replaced is None:
                rns.text = _replaceString
                replaced = True
            else:            
                rns.text = ''
                
        return paragraph

def translateppt(_source,_result):    
    prs = Presentation(_source)

    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    translator = Translator()

    text_runs = []
    translated = []

    for slide in prs.slides:                
        for shape in slide.shapes:                    
            #if not shape.has_text_frame:
            #    continue
            #if hasattr(shape,'text'):
            if shape.has_table:
                print ("translating table")
                for row in shape.table.rows:
                    for cell in row.cells:
                        text_frame = cell.text_frame                        
                        text_frame.text = translator.translate(text_frame.text, dest=LANG).text
            if shape.has_text_frame:
                print ("translating paragraph")
                for paragraph in shape.text_frame.paragraphs:      
                    if hasattr(paragraph,'text'):
                        text_runs.append(paragraph.text)                         
                        trans = translator.translate(paragraph.text, dest=LANG).text
                        paragraph = replaceText(paragraph,trans)
                        translated.append(trans)       
            
    prs.save(_result)


for dirname,subs,files in os.walk(SOURCE):
    for fname in files:           
        if 'ppt' in fname and '~$' not in fname:
            #is a power point            
            if os.path.isfile(RESULT + '/' + fname) == False:
                #is not translated already       
                translateppt(dirname + '/' + fname,RESULT + '/' + fname)
